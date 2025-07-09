import os
import fitz  # PyMuPDF
from pptx import Presentation
import requests

# === 1. PDF TEXT EXTRACTION ===
def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text.strip()

# === 2. PPTX TEXT EXTRACTION ===
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text.strip()

# === 3. CALL MISTRAL VIA OLLAMA ===
def ask_mistral(prompt):
    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": "mistral", "prompt": prompt, "stream": False}  #or mistral
        )
        return response.json()["response"].strip()
    except Exception as e:
        return f"[ERROR calling Mistral] {e}"

# === 4. CHUNK TEXT TO AVOID LLM OVERFLOW ===
def chunk_text(text, max_words=500):
    words = text.split()
    return [' '.join(words[i:i + max_words]) for i in range(0, len(words), max_words)]

# === 5. PROCESS ALL FILES ===
def process_folder(folder_path):
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)

        if file_name.endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
        elif file_name.endswith(".pptx"):
            text = extract_text_from_pptx(file_path)
        else:
            print(f"Skipping unsupported file: {file_name}")
            continue

        print(f"\n📂 Processing: {file_name}")

        chunks = chunk_text(text)
        full_summary = ""
        full_questions = ""

        for i, chunk in enumerate(chunks):
            print(f"  ➤ Chunk {i+1}...")

            summary_prompt = f"Summarize the following content to help a student revise:\n\n{chunk}"
            question_prompt = f"Generate 3 exam-style questions based on this content:\n\n{chunk}"

            summary = ask_mistral(summary_prompt)
            questions = ask_mistral(question_prompt)

            full_summary += f"\n\n[Chunk {i+1} Summary]\n" + summary
            full_questions += f"\n\n[Chunk {i+1} Questions]\n" + questions

        # Save results
        base_name = os.path.splitext(file_name)[0]
        with open(f"{base_name}_summary.txt", "w", encoding="utf-8") as f:
            f.write(full_summary.strip())
        with open(f"{base_name}_questions.txt", "w", encoding="utf-8") as f:
            f.write(full_questions.strip())

        print(f"✅ Saved: {base_name}_summary.txt and {base_name}_questions.txt")

# === 6. Run on Your Folder ===
if __name__ == "__main__":
    folder = "Input"  # Replace with your folder name
    process_folder(folder)
